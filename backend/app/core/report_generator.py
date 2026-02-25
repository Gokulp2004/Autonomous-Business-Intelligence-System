"""
report_generator.py - Automated PDF & PowerPoint Report Generation

Creates professional business intelligence reports from analysis results.

PDF: Uses fpdf2 — a pure-Python PDF library.
PPT: Uses python-pptx — creates PowerPoint presentations.

Both report types include:
  - Title page with file summary
  - Data quality overview
  - Descriptive statistics table
  - Correlation highlights
  - Trend analysis summary
  - Anomaly report
  - Forecast summary (if available)
"""

import os
import math
import datetime
from typing import Optional

import numpy as np
import pandas as pd
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

from ..config import settings


# ══════════════════════════════════════════════════════════════════
#  Helpers
# ══════════════════════════════════════════════════════════════════

def _s(v, decimals=2):
    """Safely format a value for display."""
    if v is None:
        return "-"
    if isinstance(v, (float, np.floating)):
        f = float(v)
        if math.isnan(f) or math.isinf(f):
            return "-"
        return f"{f:,.{decimals}f}"
    if isinstance(v, (int, np.integer)):
        return f"{int(v):,}"
    return str(v)


def _trunc(s, maxlen=40):
    """Truncate string for table cells."""
    s = str(s)
    return s[:maxlen] + "..." if len(s) > maxlen else s


def _ascii_safe(text: str) -> str:
    """Replace non-ASCII characters with safe ASCII equivalents for PDF."""
    replacements = {
        '\u2022': '-',   # bullet •
        '\u2192': '->',  # right arrow →
        '\u2190': '<-',  # left arrow ←
        '\u2014': '--',  # em dash —
        '\u2013': '-',   # en dash –
        '\u2018': "'",   # left single quote '
        '\u2019': "'",   # right single quote '
        '\u201c': '"',   # left double quote "
        '\u201d': '"',   # right double quote "
        '\u2026': '...', # ellipsis …
        '\u00d7': 'x',   # multiplication sign ×
        '\u2264': '<=',  # ≤
        '\u2265': '>=',  # ≥
        '\u00b1': '+/-', # ±
        '\u221e': 'inf', # ∞
    }
    for char, repl in replacements.items():
        text = text.replace(char, repl)
    # Strip any remaining non-latin1 chars
    return text.encode('latin-1', errors='replace').decode('latin-1')


# ══════════════════════════════════════════════════════════════════
#  PDF REPORT
# ══════════════════════════════════════════════════════════════════

class BIReport(FPDF):
    """Custom PDF class with header/footer."""

    def __init__(self, report_id: str):
        super().__init__()
        self.report_id = report_id
        self.set_auto_page_break(auto=True, margin=20)

    def header(self):
        if self.page_no() > 1:
            self.set_font("Helvetica", "B", 9)
            self.set_text_color(100, 100, 100)
            self.cell(0, 8, f"AutoBI Report - {self.report_id}", align="L")
            self.ln(4)
            self.set_draw_color(99, 102, 241)
            self.set_line_width(0.5)
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

    def section_title(self, title: str):
        self.ln(4)
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(30, 30, 30)
        self.cell(0, 10, _ascii_safe(title), new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(99, 102, 241)
        self.set_line_width(0.6)
        self.line(10, self.get_y(), 80, self.get_y())
        self.ln(4)

    def sub_title(self, title: str):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(60, 60, 60)
        self.cell(0, 7, _ascii_safe(title), new_x="LMARGIN", new_y="NEXT")
        self.ln(2)

    def body_text(self, text: str):
        self.set_font("Helvetica", "", 9)
        self.set_text_color(50, 50, 50)
        self.multi_cell(0, 5, _ascii_safe(text))
        self.ln(2)

    def kpi_row(self, items: list):
        """Render a row of KPI boxes. items = [(label, value), ...]"""
        col_w = 190 / len(items)
        y_start = self.get_y()
        self.set_font("Helvetica", "", 8)
        self.set_text_color(100, 100, 100)
        for label, value in items:
            x = self.get_x()
            # Box
            self.set_fill_color(245, 245, 255)
            self.rect(x, y_start, col_w - 2, 18, style="F")
            # Label
            self.set_xy(x + 2, y_start + 2)
            self.set_font("Helvetica", "", 7)
            self.set_text_color(120, 120, 120)
            self.cell(col_w - 4, 4, _ascii_safe(str(label).upper()), new_x="LEFT")
            # Value
            self.set_xy(x + 2, y_start + 7)
            self.set_font("Helvetica", "B", 12)
            self.set_text_color(30, 30, 30)
            self.cell(col_w - 4, 8, _ascii_safe(str(value)), new_x="LEFT")
            self.set_xy(x + col_w, y_start)
        self.set_y(y_start + 22)

    def add_table(self, headers: list, rows: list, col_widths: list = None):
        """Render a data table."""
        if not rows:
            return
        n = len(headers)
        if col_widths is None:
            col_widths = [190 / n] * n

        # Header
        self.set_font("Helvetica", "B", 8)
        self.set_fill_color(99, 102, 241)
        self.set_text_color(255, 255, 255)
        for i, h in enumerate(headers):
            self.cell(col_widths[i], 7, _ascii_safe(_trunc(str(h), 20)), border=1, fill=True, align="C")
        self.ln()

        # Rows
        self.set_font("Helvetica", "", 7.5)
        self.set_text_color(40, 40, 40)
        for ri, row in enumerate(rows):
            if self.get_y() > 265:
                self.add_page()
            fill = ri % 2 == 0
            if fill:
                self.set_fill_color(248, 248, 252)
            for i, val in enumerate(row):
                align = "R" if i > 0 else "L"
                self.cell(col_widths[i], 6, _ascii_safe(_trunc(_s(val) if isinstance(val, (int, float, np.integer, np.floating)) else str(val), 22)), border=1, fill=fill, align=align)
            self.ln()
        self.ln(3)


def generate_pdf_report(
    file_id: str,
    results: dict,
    df: pd.DataFrame,
) -> str:
    """
    Generate a comprehensive PDF report.

    Returns the file path to the generated PDF.
    """
    pdf = BIReport(file_id)
    pdf.alias_nb_pages()

    analysis = results.get("analysis", {})
    cleaning = results.get("cleaning", {})
    anomalies = results.get("anomalies", {})
    forecasts = results.get("forecasts")
    summary = analysis.get("summary", {})

    # ── Title Page ────────────────────────────────────────────
    pdf.add_page()
    pdf.ln(40)
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(99, 102, 241)
    pdf.cell(0, 15, "Business Intelligence Report", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, f"Dataset: {file_id}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 8, f"Generated: {datetime.datetime.now().strftime('%B %d, %Y at %I:%M %p')}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(15)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 6, f"{_s(summary.get('total_rows', len(df)))} rows  |  {_s(summary.get('total_columns', len(df.columns)))} columns  |  {_s(summary.get('numeric_columns', 0))} numeric  |  {_s(summary.get('categorical_columns', 0))} categorical", align="C", new_x="LMARGIN", new_y="NEXT")

    # ── Data Quality Overview ─────────────────────────────────
    pdf.add_page()
    pdf.section_title("1. Data Quality Overview")

    clean_sum = cleaning.get("summary", {})
    total_cells = summary.get("total_rows", 0) * summary.get("total_columns", 0)
    missing = int(df.isnull().sum().sum()) if total_cells > 0 else 0
    quality = round((1 - missing / total_cells) * 100, 1) if total_cells > 0 else 100.0

    pdf.kpi_row([
        ("Data Quality", f"{quality}%"),
        ("Rows After Cleaning", _s(clean_sum.get("rows_after", len(df)))),
        ("Duplicates Removed", _s(clean_sum.get("duplicates_removed", 0))),
        ("Missing Values Fixed", _s(clean_sum.get("missing_after", 0))),
    ])

    actions = cleaning.get("actions", [])
    if actions:
        pdf.sub_title("Cleaning Steps Performed")
        for act in actions[:12]:
            step = act.get("step", "")
            detail = act.get("detail", "")
            pdf.set_font("Helvetica", "B", 8)
            pdf.set_text_color(99, 102, 241)
            pdf.cell(4, 5, "-")
            pdf.set_font("Helvetica", "B", 8)
            pdf.set_text_color(40, 40, 40)
            pdf.cell(40, 5, _ascii_safe(step))
            pdf.set_font("Helvetica", "", 8)
            pdf.set_text_color(80, 80, 80)
            pdf.cell(0, 5, _ascii_safe(_trunc(detail, 90)), new_x="LMARGIN", new_y="NEXT")

    # ── Descriptive Statistics ────────────────────────────────
    pdf.add_page()
    pdf.section_title("2. Descriptive Statistics")

    desc = analysis.get("descriptive_stats", {})
    if desc:
        cols = list(desc.keys())
        metrics = ["count", "mean", "std", "min", "median", "max", "cv"]
        headers = ["Metric"] + [_trunc(c, 14) for c in cols]
        widths = [24] + [(190 - 24) / len(cols)] * len(cols)
        rows = []
        for m in metrics:
            row = [m.upper()]
            for c in cols:
                row.append(desc[c].get(m))
            rows.append(row)
        pdf.add_table(headers, rows, widths)

    # ── Correlations ──────────────────────────────────────────
    strong = analysis.get("strong_correlations", [])
    if strong:
        pdf.section_title("3. Strong Correlations")
        pdf.body_text(f"Found {len(strong)} strong correlation(s) with |r| > 0.7:")
        headers = ["Column A", "Column B", "r", "Strength", "Direction"]
        rows = [[c.get("col_a"), c.get("col_b"), c.get("correlation"), c.get("strength"), c.get("direction")] for c in strong]
        pdf.add_table(headers, rows, [45, 45, 25, 40, 35])
    else:
        pdf.section_title("3. Correlations")
        pdf.body_text("No strong correlations (|r| > 0.7) found in this dataset.")

    # ── Distributions ─────────────────────────────────────────
    dists = analysis.get("distributions", {})
    if dists:
        pdf.section_title("4. Distribution Analysis")
        headers = ["Column", "Skewness", "Kurtosis", "Shape", "Normal?"]
        rows = []
        for col, d in dists.items():
            rows.append([col, d.get("skewness"), d.get("kurtosis"), d.get("shape", "").replace("_", " "), "Yes" if d.get("is_normal") else "No"])
        pdf.add_table(headers, rows, [40, 30, 30, 45, 25])

    # ── Trend Analysis ────────────────────────────────────────
    trends = analysis.get("trends", {})
    if trends:
        pdf.add_page()
        pdf.section_title("5. Trend Analysis")
        headers = ["Column", "Direction", "Slope", "R-squared", "p-value", "Significant?", "% Change"]
        rows = []
        for col, t in trends.items():
            rows.append([
                col, t.get("direction"), t.get("slope"), t.get("r_squared"),
                t.get("p_value"), "Yes" if t.get("significant") else "No",
                f"{t.get('pct_change')}%" if t.get("pct_change") is not None else "-",
            ])
        pdf.add_table(headers, rows, [28, 24, 26, 26, 26, 24, 26])

    # ── Anomalies ─────────────────────────────────────────────
    anom_summary = anomalies.get("summary", {})
    per_col = anomalies.get("per_column", {})
    pdf.section_title("6. Anomaly Detection")
    total_anom = anom_summary.get("total_anomalies", 0)
    pdf.kpi_row([
        ("Total Anomalies", _s(total_anom)),
        ("Columns Affected", _s(anom_summary.get("columns_with_anomalies", 0))),
        ("Methods Used", "Z-score + IQR + IF"),
    ])
    if per_col:
        headers = ["Column", "Z-score", "IQR", "Total"]
        rows = []
        for col, info in per_col.items():
            z_count = info.get("zscore", {}).get("count", 0) if isinstance(info.get("zscore"), dict) else 0
            i_count = info.get("iqr", {}).get("count", 0) if isinstance(info.get("iqr"), dict) else 0
            rows.append([col, z_count, i_count, info.get("total_anomalies", 0)])
        pdf.add_table(headers, rows, [50, 40, 40, 40])

    # ── Forecasts ─────────────────────────────────────────────
    if forecasts and isinstance(forecasts, list) and len(forecasts) > 0:
        pdf.add_page()
        pdf.section_title("7. Forecast Results")
        pdf.body_text(f"Forecasts generated for {len(forecasts)} column(s).")
        for fc in forecasts[:6]:
            col = fc.get("column", "Unknown")
            method = fc.get("method", "Unknown")
            predictions = fc.get("forecast", [])
            pdf.sub_title(f"{col} (method: {method})")
            if predictions:
                headers = ["Date", "Predicted", "Lower CI", "Upper CI"]
                rows = []
                for p in predictions[:10]:
                    rows.append([
                        str(p.get("date", "-"))[:10],
                        p.get("predicted"),
                        p.get("lower"),
                        p.get("upper"),
                    ])
                pdf.add_table(headers, rows, [50, 45, 45, 45])

    # ── Feature Importance ────────────────────────────────────
    fi = analysis.get("feature_importance", [])
    if fi:
        pdf.section_title("8. Feature Importance")
        headers = ["Feature", "Target", "Importance"]
        rows = [[f.get("feature"), f.get("target"), f.get("importance")] for f in fi]
        pdf.add_table(headers, rows, [60, 60, 60])

    # ── Save ──────────────────────────────────────────────────
    output_dir = settings.OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{file_id}_report.pdf"
    filepath = os.path.join(output_dir, filename)
    pdf.output(filepath)
    return filename


# ══════════════════════════════════════════════════════════════════
#  POWERPOINT REPORT
# ══════════════════════════════════════════════════════════════════

# Color constants
_INDIGO = RGBColor(99, 102, 241)
_DARK = RGBColor(30, 30, 30)
_GRAY = RGBColor(100, 100, 100)
_WHITE = RGBColor(255, 255, 255)
_LIGHT_BG = RGBColor(245, 245, 255)


def _add_title_slide(prs: Presentation, file_id: str, summary: dict):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(99, 102, 241)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Business Intelligence Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"Dataset: {file_id}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Date
    p3 = tf.add_paragraph()
    p3.text = datetime.datetime.now().strftime("%B %d, %Y")
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(200, 200, 255)
    p3.alignment = PP_ALIGN.CENTER

    # Stats bar
    stats_text = f"{_s(summary.get('total_rows', 0))} rows  |  {_s(summary.get('total_columns', 0))} columns  |  {_s(summary.get('numeric_columns', 0))} numeric"
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p4 = tf2.paragraphs[0]
    p4.text = stats_text
    p4.font.size = Pt(11)
    p4.font.color.rgb = RGBColor(200, 200, 255)
    p4.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, title: str):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _INDIGO
    p.alignment = PP_ALIGN.CENTER
    return slide


def _add_content_slide(prs: Presentation, title: str) -> object:
    """Add a slide with a title and return (slide, top_y)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = _DARK
    return slide


def _add_table_slide(prs: Presentation, title: str, headers: list, rows: list, col_widths: list = None):
    """Add a slide with a data table."""
    slide = _add_content_slide(prs, title)

    n_cols = len(headers)
    n_rows = min(len(rows), 14) + 1  # cap at 14 data rows + header
    total_w = Inches(9)
    table_h = Inches(0.35 * n_rows)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.1), total_w, table_h)
    tbl = tbl_shape.table

    # Set column widths
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(total_w.emu * w / total))

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _INDIGO
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, row in enumerate(rows[:14]):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            display = _s(val) if isinstance(val, (int, float, np.integer, np.floating)) else str(val if val is not None else "-")
            cell.text = _trunc(display, 25)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _LIGHT_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = _DARK
                paragraph.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def _add_kpi_slide(prs: Presentation, title: str, kpis: list):
    """Add a slide with KPI boxes. kpis = [(label, value, color_hex), ...]"""
    slide = _add_content_slide(prs, title)
    n = len(kpis)
    box_w = Inches(8.5 / n)
    box_h = Inches(1.2)
    start_x = Inches(0.5)
    start_y = Inches(1.5)

    for i, (label, value, *rest) in enumerate(kpis):
        color = rest[0] if rest else "6366f1"
        x = start_x + box_w * i + Inches(0.1) * i

        # Box background
        shape = slide.shapes.add_shape(
            1, x, start_y, box_w - Inches(0.1), box_h
        )
        shape.fill.solid()
        r, g, b = int(color[:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

        # Label
        txB = slide.shapes.add_textbox(x + Inches(0.15), start_y + Inches(0.1), box_w - Inches(0.3), Inches(0.3))
        p = txB.text_frame.paragraphs[0]
        p.text = str(label).upper()
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(min(r + 80, 255), min(g + 80, 255), min(b + 80, 255))

        # Value
        txV = slide.shapes.add_textbox(x + Inches(0.15), start_y + Inches(0.4), box_w - Inches(0.3), Inches(0.6))
        pv = txV.text_frame.paragraphs[0]
        pv.text = str(value)
        pv.font.size = Pt(24)
        pv.font.bold = True
        pv.font.color.rgb = _WHITE

    return slide


def _add_bullet_slide(prs: Presentation, title: str, items: list):
    """Add a slide with bullet points. items = [(bold_text, detail_text), ...]."""
    slide = _add_content_slide(prs, title)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items[:12]):
        if isinstance(item, tuple):
            bold_part, detail = item
        else:
            bold_part, detail = str(item), ""

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)

        run1 = p.add_run()
        run1.text = f"\u2022 {bold_part}"
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = _DARK

        if detail:
            run2 = p.add_run()
            run2.text = f"  -  {detail}"
            run2.font.size = Pt(10)
            run2.font.color.rgb = _GRAY

    return slide


def generate_ppt_report(
    file_id: str,
    results: dict,
    df: pd.DataFrame,
) -> str:
    """
    Generate a professional PowerPoint report.

    Returns the filename of the generated PPTX.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    analysis = results.get("analysis", {})
    cleaning = results.get("cleaning", {})
    anomalies = results.get("anomalies", {})
    forecasts = results.get("forecasts")
    summary = analysis.get("summary", {})
    clean_sum = cleaning.get("summary", {})

    # 1. Title Slide
    _add_title_slide(prs, file_id, summary)

    # 2. Data Quality KPIs
    total_cells = summary.get("total_rows", 0) * summary.get("total_columns", 0)
    missing = int(df.isnull().sum().sum()) if total_cells > 0 else 0
    quality = round((1 - missing / total_cells) * 100, 1) if total_cells > 0 else 100.0
    _add_kpi_slide(prs, "Data Quality Overview", [
        ("Quality Score", f"{quality}%", "10b981"),
        ("Total Rows", _s(summary.get("total_rows", len(df))), "6366f1"),
        ("Columns", _s(summary.get("total_columns", len(df.columns))), "3b82f6"),
        ("Duplicates Removed", _s(clean_sum.get("duplicates_removed", 0)), "f59e0b"),
    ])

    # 3. Descriptive Statistics
    desc = analysis.get("descriptive_stats", {})
    if desc:
        cols = list(desc.keys())
        metrics = ["count", "mean", "std", "min", "median", "max"]
        headers = ["Metric"] + cols
        rows = []
        for m in metrics:
            row = [m.upper()]
            for c in cols:
                row.append(desc[c].get(m))
            rows.append(row)
        _add_table_slide(prs, "Descriptive Statistics", headers, rows)

    # 4. Correlations
    strong = analysis.get("strong_correlations", [])
    if strong:
        headers = ["Column A", "Column B", "r", "Strength", "Direction"]
        rows = [[c.get("col_a"), c.get("col_b"), c.get("correlation"), c.get("strength"), c.get("direction")] for c in strong]
        _add_table_slide(prs, "Strong Correlations (|r| > 0.7)", headers, rows)

    # 5. Distributions
    dists = analysis.get("distributions", {})
    if dists:
        headers = ["Column", "Skewness", "Kurtosis", "Shape", "Normal?"]
        rows = [[col, d.get("skewness"), d.get("kurtosis"), d.get("shape", "").replace("_", " "), "Yes" if d.get("is_normal") else "No"] for col, d in dists.items()]
        _add_table_slide(prs, "Distribution Analysis", headers, rows)

    # 6. Trends
    trends = analysis.get("trends", {})
    if trends:
        items = []
        for col, t in trends.items():
            direction = t.get("direction", "flat")
            pct = f"{t.get('pct_change')}% change" if t.get("pct_change") is not None else ""
            sig = "significant" if t.get("significant") else "not significant"
            items.append((f"{col}: {direction}", f"R-sq={t.get('r_squared')}, {sig}, {pct}"))
        _add_bullet_slide(prs, "Trend Analysis", items)

    # 7. Anomalies
    anom_summary = anomalies.get("summary", {})
    per_col = anomalies.get("per_column", {})
    total_anom = anom_summary.get("total_anomalies", 0)
    _add_kpi_slide(prs, "Anomaly Detection", [
        ("Total Anomalies", _s(total_anom), "ef4444" if total_anom > 0 else "10b981"),
        ("Columns Affected", _s(anom_summary.get("columns_with_anomalies", 0)), "f59e0b"),
        ("Detection Methods", "3", "6366f1"),
    ])

    if per_col:
        headers = ["Column", "Z-score", "IQR", "Total"]
        rows = []
        for col, info in per_col.items():
            z_count = info.get("zscore", {}).get("count", 0) if isinstance(info.get("zscore"), dict) else 0
            i_count = info.get("iqr", {}).get("count", 0) if isinstance(info.get("iqr"), dict) else 0
            rows.append([col, z_count, i_count, info.get("total_anomalies", 0)])
        _add_table_slide(prs, "Anomalies by Column", headers, rows)

    # 8. Forecasts
    if forecasts and isinstance(forecasts, list):
        items = []
        for fc in forecasts[:8]:
            col = fc.get("column", "Unknown")
            method = fc.get("method", "Unknown")
            preds = fc.get("forecast", [])
            n = len(preds)
            items.append((f"{col} ({method})", f"{n} period forecast"))
        if items:
            _add_bullet_slide(prs, "Forecast Summary", items)

    # 9. Feature Importance
    fi = analysis.get("feature_importance", [])
    if fi:
        headers = ["Feature", "Target", "Importance"]
        rows = [[f.get("feature"), f.get("target"), f.get("importance")] for f in fi]
        _add_table_slide(prs, "Feature Importance", headers, rows)

    # Save
    output_dir = settings.OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{file_id}_report.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filename
